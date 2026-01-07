"""
**PptxSlideGenerator Module**
------------------------------

The `PptxSlideGenerator` module provides PowerPoint slide generation based on lesson objectives,
readings, and prior lesson presentations. This is the concrete implementation of the BaseSlideGenerator
for PowerPoint (.pptx) output format.

Key Functionalities
~~~~~~~~~~~~~~~~~~~

1. **Automated PowerPoint Slide Generation**:
   - Generates a PowerPoint presentation for each lesson, incorporating:
     - A title slide with consistent author and institution information
     - "Where We Came From" and "Where We Are Going" slides
     - Lesson objectives with highlighted action verbs
     - Discussion questions and in-class exercises
     - Summary slides with key takeaways

2. **Previous Lesson Integration**:
   - Retrieves and references prior lesson presentations to maintain consistent formatting and flow
   - Preserves author and institution information across presentations

3. **PowerPoint-Specific Features**:
   - Uses python-pptx library for PowerPoint generation
   - Validates slide structure and content
   - Supports custom layouts and themes

Status
~~~~~~
**PLACEHOLDER FOR FUTURE IMPLEMENTATION**

This module is currently a skeleton/placeholder. The implementation will include:
- PowerPoint-specific prompts (no LaTeX syntax)
- PptxSlides schema for structured output
- python-pptx integration for .pptx file creation
- PowerPoint-specific validation

Usage (Planned)
~~~~~~~~~~~~~~~

1. **Initialize PptxSlideGenerator**:
   ```python
   generator = PptxSlideGenerator(
       lesson_no=10,
       llm=llm,
       course_name="Political Science",
       lesson_loader=lesson_loader,
       output_dir=output_dir
   )
   ```

2. **Generate Slides**:
   ```python
   guidance = "Focus on comparing democratic and authoritarian systems"
   slides = generator.generate_slides(specific_guidance=guidance)
   ```

3. **Save the Slides**:
   ```python
   generator.save_slides(slides)
   ```
"""
# %%
import logging
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Union

from langchain_core.messages import SystemMessage
from langchain_core.prompts import (ChatPromptTemplate,
                                    HumanMessagePromptTemplate)

from class_factory.beamer_bot.base_slide_generator import BaseSlideGenerator
from class_factory.beamer_bot.pptx_prompts import (pptx_human_prompt,
                                                   pptx_system_prompt)
from class_factory.beamer_bot.pptx_slides import PptxSlides
from class_factory.utils.load_documents import LessonLoader


class PptxSlideGenerator(BaseSlideGenerator):
    """
    A class to generate PowerPoint slides for a specified lesson using a language model (LLM).

    **NOTE: This is currently a placeholder implementation.**

    PptxSlideGenerator will automate the slide generation process, creating structured presentations
    based on lesson readings, objectives, and content from prior presentations when available.

    Attributes:
        lesson_no (int): Lesson number for which to generate slides.
        llm: Language model instance for generating slides.
        course_name (str): Name of the course for slide context.
        lesson_loader (LessonLoader): Loader for accessing lesson readings and objectives.
        output_dir (Path): Directory to save the generated PowerPoint slides.
        slide_dir (Path): Directory containing existing PowerPoint slides.
        llm_response: Stores the generated response from the LLM.
        prompt: Generated prompt for the LLM.
        pptx_output (Path): Output path for the generated .pptx file.

    Methods:
        generate_slides(specific_guidance: str = None) -> Any:
            Generates PowerPoint slides for the specified lesson.

        save_slides(content: Any) -> None:
            Saves the generated PowerPoint content to a .pptx file.
    """

    def __init__(self, lesson_no: int, llm, course_name: str, lesson_loader: LessonLoader,
                 output_dir: Union[Path, str] = None, verbose: bool = False,
                 slide_dir: Union[Path, str] = None, lesson_objectives: dict = None):
        """
        Initialize the PowerPoint slide generator.

        Args:
            lesson_no (int): Lesson number for which to generate slides
            llm: Language model instance for content generation
            course_name (str): Name of the course
            lesson_loader (LessonLoader): Loader for lesson materials
            output_dir (Union[Path, str], optional): Output directory for generated slides
            verbose (bool, optional): Enable verbose logging. Defaults to False.
            slide_dir (Union[Path, str], optional): Directory with existing slides
            lesson_objectives (dict, optional): User-provided lesson objectives
        """
        super().__init__(
            lesson_no=lesson_no,
            llm=llm,
            course_name=course_name,
            lesson_loader=lesson_loader,
            output_dir=output_dir,
            verbose=verbose,
            slide_dir=slide_dir,
            lesson_objectives=lesson_objectives
        )

        # Generate the prompt (to be implemented)
        self.prompt = self._generate_prompt()

        # Use LLM with structured output for PowerPoint slides
        self.chain = self.prompt | self.llm.with_structured_output(PptxSlides)

        # Set output file path
        self.pptx_output = self.output_dir / f'L{self.lesson_no}.pptx'

    def _load_prior_lesson(self) -> str:
        """
        Load the previous lesson's PowerPoint presentation.

        Returns:
            str: Text representation of the prior lesson's PowerPoint presentation

        Note:
            This attempts to extract text from a .pptx file or falls back to .tex files.
        """
        try:
            # Try to find a prior .pptx file
            prior_pptx = self._find_prior_lesson(self.lesson_no, max_attempts=3, file_extension=".pptx")

            # Extract text from PowerPoint
            from pptx import Presentation
            prs = Presentation(str(prior_pptx))

            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"Slide {i}:\n"

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"

                slides_text.append(slide_text)

            return "\n\n".join(slides_text)

        except FileNotFoundError:
            self.logger.info("No prior PowerPoint found, trying LaTeX files...")
            try:
                # Fall back to LaTeX files if PowerPoint not found
                beamer_example = self.lesson_loader.find_prior_beamer_presentation(self.lesson_no)
                return self.lesson_loader.load_beamer_presentation(beamer_example)
            except Exception as e:
                self.logger.warning(f"Could not load prior presentation: {e}")
                return "Not Provided"
        except ImportError:
            self.logger.warning("python-pptx not installed, cannot read prior PowerPoint files")
            return "Not Provided"
        except Exception as e:
            self.logger.warning(f"Error loading prior lesson: {e}")
            return "Not Provided"

    def _generate_prompt(self, human_prompt: str = None):
        """
        Generates a detailed prompt for the LLM to guide PowerPoint slide creation.

        Args:
            human_prompt (str, optional): Override default human prompt template

        Returns:
            ChatPromptTemplate: The constructed prompt for the LLM.
        """
        prompt = ChatPromptTemplate.from_messages(
            [
                SystemMessage(
                    content=(
                        pptx_system_prompt.format(
                            course_name=self.course_name)
                    )
                ),
                HumanMessagePromptTemplate.from_template(
                    pptx_human_prompt if not human_prompt else human_prompt),
            ]
        )

        return prompt

    def generate_slides(self, specific_guidance: str = None, lesson_objectives: dict = None, **kwargs) -> PptxSlides:
        """
        Generate PowerPoint slides for the lesson using the language model.

        Args:
            specific_guidance (str, optional): Custom instructions for slide content and structure
            lesson_objectives (dict, optional): Override default objectives with custom ones
            **kwargs: Additional PowerPoint-specific arguments

        Returns:
            PptxSlides: Generated slide data as PptxSlides object

        Raises:
            ValueError: If validation fails after maximum retry attempts

        Note:
            This implementation generates slides using LLM with PowerPoint-specific prompts.
            The slides use plain text with markdown-like formatting (no LaTeX).
        """
        # Load objectives (last, current, next), readings, and previous lesson slides
        self.user_objectives = (
            self.set_user_objectives(lesson_objectives, range(self.lesson_no, self.lesson_no + 1))
            if lesson_objectives else {}
        )
        objectives_text = "\n\n".join([
            self._get_lesson_objectives(lesson)
            for lesson in range(self.lesson_no - 1, self.lesson_no + 2)
        ])
        combined_readings_text = self.readings

        if self.lesson_loader.slide_dir:
            try:
                prior_lesson_content = self._load_prior_lesson()
            except Exception as e:
                self.logger.warning(f"Could not load prior lesson: {e}")
                prior_lesson_content = "Not Provided"
        else:
            prior_lesson_content = "Not Provided"
            self.logger.warning(
                "No slide_dir provided. Prior slides will not be referenced during generation."
            )

        self.logger.info("Generating PowerPoint slides...")

        # Generate PowerPoint slides via the chain
        additional_guidance = ""
        retries, MAX_RETRIES = 0, 3
        valid = False

        while not valid and retries < MAX_RETRIES:
            try:
                # LLM returns PptxSlides object
                slides_data = self.chain.invoke({
                    "objectives": objectives_text,
                    "information": combined_readings_text,
                    "last_presentation": prior_lesson_content,
                    "lesson_no": self.lesson_no,
                    'specific_guidance': specific_guidance if specific_guidance else "Not provided.",
                    "additional_guidance": additional_guidance
                })

                self.llm_response = slides_data

                # Basic validation - check if we have slides
                if not slides_data.slides or len(slides_data.slides) == 0:
                    raise ValueError("No slides were generated")

                # Validate slide structure
                val_response = self._validate_llm_response(
                    generated_slides=slides_data,
                    objectives=objectives_text,
                    readings=combined_readings_text,
                    last_presentation=prior_lesson_content,
                    prompt_specific_guidance=specific_guidance if specific_guidance else "Not provided.",
                    task_schema=PptxSlides.model_json_schema()
                )

                self.validator.validation_result = val_response
                self.validator.logger.info(f"Validation output: {val_response}")

                # Check validation status
                if int(val_response.get('status', 0)) != 1:
                    retries += 1
                    additional_guidance = val_response.get("additional_guidance", "")
                    self.validator.logger.warning(
                        f"Response validation failed on attempt {retries}. "
                        f"Guidance for improvement: {additional_guidance}"
                    )
                    continue  # Retry LLM generation

                valid = True
                return slides_data

            except Exception as e:
                retries += 1
                self.logger.error(f"Error generating slides (attempt {retries}): {e}")
                if retries >= MAX_RETRIES:
                    raise

        # Handle validation failure after max retries
        if not valid:
            raise ValueError(
                "Validation failed after max retries. Ensure correct prompt and input data. "
                "Consider trying a different LLM."
            )

    def _validate_llm_response(self, generated_slides: PptxSlides, objectives: str,
                               readings: str, last_presentation: str,
                               prompt_specific_guidance: str = "",
                               additional_guidance: str = "",
                               task_schema=None) -> Dict[str, Any]:
        """
        Validates the generated PowerPoint slides for content quality and structure.

        Args:
            generated_slides (PptxSlides): PowerPoint slides generated by the LLM
            objectives (str): Formatted string of lesson objectives for validation
            readings (str): Formatted string of lesson readings for content verification
            last_presentation (str): Content from prior lesson's presentation for format consistency
            prompt_specific_guidance (str, optional): Custom guidance provided during generation
            additional_guidance (str, optional): Supplementary guidance for validation refinement
            task_schema: JSON schema for the PptxSlides model

        Returns:
            Dict[str, Any]: Validation results containing:
                - status (int): 1 for pass, 0 for fail
                - evaluation_score (float): Quality score (0-10)
                - additional_guidance (str): Suggestions for improvement if validation fails
                - rationale (str): Explanation of the validation result
        """
        validation_prompt = self.prompt.format(
            objectives=objectives,
            information=readings,
            last_presentation=last_presentation,
            lesson_no=self.lesson_no,
            additional_guidance=additional_guidance,
            specific_guidance=prompt_specific_guidance
        )

        val_response = self.validator.validate(
            task_description=validation_prompt,
            generated_response=generated_slides,
            task_schema=task_schema,
            specific_guidance=prompt_specific_guidance + (
                "\n" + additional_guidance if additional_guidance else ""
            )
        )
        return val_response

    def save_slides(self, content: PptxSlides, output_dir: Union[Path, str] = None) -> None:
        """
        Save the generated PowerPoint slides to a .pptx file.

        Args:
            content (PptxSlides): The slide content to save
            output_dir (Union[Path, str], optional): Override default output directory

        Raises:
            ImportError: If python-pptx is not installed

        Note:
            This method uses python-pptx to create the PowerPoint file.
            Install with: pip install python-pptx
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError(
                "python-pptx is required to save PowerPoint files. "
                "Install it with: pip install python-pptx"
            )

        # Always set date to current date
        content.date = datetime.now().strftime("%B %Y")

        # Determine output path
        if output_dir:
            pptx_file = Path(output_dir) / f'L{self.lesson_no}.pptx'
        else:
            pptx_file = self.pptx_output

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Process each slide
        for slide_data in content.slides:

            # Determine layout based on slide type
            if slide_data.slide_type == "title":
                slide_layout = prs.slide_layouts[0]  # Title Slide
                slide = prs.slides.add_slide(slide_layout)

                # Set title
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.title

                # Set subtitle with author, institute, and date from metadata
                if len(slide.shapes) > 1:
                    subtitle_shape = slide.shapes[1]
                    subtitle_parts = []

                    if content.author:
                        subtitle_parts.append(content.author)
                    if content.institute:
                        subtitle_parts.append(content.institute)
                    if content.date:
                        subtitle_parts.append(content.date)

                    if subtitle_parts:
                        subtitle_shape.text = "\n".join(subtitle_parts)
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)

                # Set title
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.title

                # Add content
                if len(slide.shapes) > 1:
                    # Get the content placeholder
                    content_shape = slide.shapes[1]
                    text_frame = content_shape.text_frame
                    text_frame.word_wrap = True

                    # Check if content has markdown-style bullets (starts with * or -)
                    # Or multiple distinct lines that should be bullets
                    if slide_data.content and not slide_data.bullet_points:
                        content = slide_data.content.strip()
                        lines = content.split('\n')

                        # Check for markdown-style bullets (lines starting with *)
                        bullet_lines = [line.strip() for line in lines if line.strip().startswith('*')]

                        # If we have 2+ bullet lines, convert them
                        if len(bullet_lines) >= 2:
                            slide_data.bullet_points = [line.lstrip('*').strip() for line in bullet_lines]
                            # Keep non-bullet content as content
                            non_bullet_lines = [line.strip() for line in lines if not line.strip().startswith('*') and line.strip()]
                            slide_data.content = '\n'.join(non_bullet_lines) if non_bullet_lines else ''

                        # Alternatively, check if we have 3+ separate lines (separated by double newlines or single lines)
                        # that look like they should be bullets (short statements)
                        elif '\n\n' in content:
                            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
                            # If we have 3+ short paragraphs (< 100 chars each), convert to bullets
                            if len(paragraphs) >= 3 and all(len(p) < 100 for p in paragraphs):
                                slide_data.bullet_points = paragraphs
                                slide_data.content = ''

                    # Add bullet points if present
                    if slide_data.bullet_points:
                        # Don't clear - preserve bullet formatting from layout
                        # Just clear the text from the first paragraph
                        if text_frame.paragraphs:
                            text_frame.paragraphs[0].text = ""

                        # Get bullet format from first paragraph to copy to others
                        first_p = text_frame.paragraphs[0]

                        for i, bullet in enumerate(slide_data.bullet_points):
                            if i == 0:
                                # Use the existing first paragraph (preserves bullet formatting)
                                p = first_p
                            else:
                                # Add new paragraph and copy bullet formatting from first
                                p = text_frame.add_paragraph()
                                p.level = first_p.level
                                # Copy bullet format
                                if first_p.font.name:
                                    p.font.name = first_p.font.name

                            p.text = self._strip_markdown(bullet)
                            p.font.size = Pt(20)

                    # Add main content if present (no bullets)
                    elif slide_data.content:
                        text_frame.clear()
                        p = text_frame.paragraphs[0]
                        p.text = self._strip_markdown(slide_data.content)
                        p.level = 0
                        p.font.size = Pt(18)

            # Add speaker notes if present
            if slide_data.notes:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = slide_data.notes

        # Save the presentation
        prs.save(str(pptx_file))
        self.logger.info(f"PowerPoint slides saved to {pptx_file}")

    def _strip_markdown(self, text: str) -> str:
        """
        Strip markdown-like formatting from text for PowerPoint.

        Args:
            text: Text with markdown formatting

        Returns:
            Plain text without markdown symbols
        """
        # Remove bold markers
        text = text.replace("**", "")
        # Remove italic markers
        text = text.replace("*", "")
        return text


# %%
if __name__ == "__main__":
    import os

    import yaml
    from dotenv import load_dotenv
    from langchain_community.llms import Ollama
    from langchain_google_genai import ChatGoogleGenerativeAI
    from langchain_openai import ChatOpenAI
    from pyprojroot.here import here

    from class_factory.utils.tools import reset_loggers

    wd = here()
    load_dotenv()

    user_home = Path.home()
    reset_loggers(log_level=logging.INFO)

    OPENAI_KEY = os.getenv('openai_key')
    OPENAI_ORG = os.getenv('openai_org')
    GEMINI_KEY = os.getenv('gemini_api_key')

    # Path definitions
    with open("class_config.yaml", "r") as file:
        config = yaml.safe_load(file)

    # class_config = config['PS491']
    class_config = config['PS300']

    slide_dir = user_home / class_config['slideDir']
    syllabus_path = user_home / class_config['syllabus_path']
    readingsDir = user_home / class_config['reading_dir']
    is_tabular_syllabus = class_config['is_tabular_syllabus']

    # llm = ChatOpenAI(
    #     model="gpt-4o-mini",
    #     temperature=0.4,
    #     max_tokens=None,
    #     timeout=None,
    #     max_retries=2,
    #     api_key=OPENAI_KEY,
    #     organization=OPENAI_ORG,
    # )

    llm = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        temperature=0.4,
        max_tokens=None,
        timeout=None,
        max_retries=2,
        api_key=GEMINI_KEY
    )

    lsn = 3

    # llm = Ollama(
    #     model="llama3.1",
    #     temperature=0.2
    # )

    specific_guidance = """
    The objectives slide should include an objective titled "have tons of fun"
    """

    loader = LessonLoader(syllabus_path=syllabus_path,
                          reading_dir=readingsDir,
                          slide_dir=slide_dir,
                          tabular_syllabus=is_tabular_syllabus
                          )

    # Initialize the BeamerBot (defaults to LaTeX)
    beamer_bot = PptxSlideGenerator(
        lesson_no=2,
        lesson_loader=loader,
        llm=llm,
        course_name="Political Science Research Methods",
        verbose=True
    )

    # Generate slides
    slides = beamer_bot.generate_slides()

    print(slides)
    # Save the generated PowerPoint slides
    beamer_bot.save_slides(slides)


# %%
