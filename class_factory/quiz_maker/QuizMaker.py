"""
**QuizMaker Module**
---------------------

The `QuizMaker` module offers a comprehensive framework for generating, distributing, and analyzing quiz questions based on lesson content and objectives. At its core, the `QuizMaker` class uses a language model (LLM) to create targeted quiz questions, ensuring these questions are relevant to the course material. This class also provides utilities for similarity checking, interactive quiz launches, and detailed results assessment.

Key Functionalities:
~~~~~~~~~~~~~~~~~~~~

1. **Quiz Generation**:
   - Automatically generates quiz questions from lesson objectives and readings.
   - Customizable difficulty level and quiz content based on user-provided or auto-extracted lesson objectives.

2. **Similarity Checking**:
   - Detects overlap with previous quizzes to prevent question duplication.
   - Uses sentence embeddings to flag and remove questions too similar to prior quizzes.

3. **Validation and Formatting**:
   - Validates generated questions to ensure proper format and structure.
   - Corrects answer formatting to meet quiz standards (e.g., answers in 'A', 'B', 'C', 'D').

4. **Saving Quizzes**:
   - Exports quizzes as Excel files or PowerPoint presentations.
   - Customizes PowerPoint presentations using templates for polished quiz slides.

5. **Interactive Quiz Launch**:
   - Launches an interactive Gradio-based web quiz for real-time participation.
   - Supports QR code access and real-time result saving.

6. **Results Assessment**:
   - Analyzes and visualizes quiz results stored in CSV files.
   - Generates summary statistics, HTML reports, and dashboards for insights into quiz performance.

Dependencies
~~~~~~~~~~~~~

This module requires:

- `langchain_core`: For LLM interaction and prompt handling.
- `sentence_transformers`: For semantic similarity detection in quiz questions.
- `pptx`: For PowerPoint presentation generation.
- `pandas`: For data handling and result assessment.
- `torch`: For managing device selection and embedding models.
- `gradio`: For interactive quiz interfaces.
- Custom utilities for document loading, response parsing, logging, and retry decorators.

Usage Overview
~~~~~~~~~~~~~~

1. **Initialize QuizMaker**:
   - Instantiate `QuizMaker` with required paths, lesson loader, and LLM.

2. **Generate a Quiz**:
   - Call `make_a_quiz()` to create quiz questions based on lesson materials, with automatic similarity checking.

3. **Save the Quiz**:
   - Use `save_quiz()` to save the quiz as an Excel file or `save_quiz_to_ppt()` to export to PowerPoint.

4. **Launch an Interactive Quiz**:
   - Use `launch_interactive_quiz()` to start a web-based quiz, with options for real-time participation and result saving.

5. **Assess Quiz Results**:
   - Analyze saved quiz responses with `assess_quiz_results()`, generating summary statistics, reports, and visualizations.

Example
~~~~~~~~

.. code-block:: python

    from pathlib import Path
    from class_factory.quiz_maker.QuizMaker import QuizMaker
    from class_factory.utils.load_documents import LessonLoader
    from langchain_openai import ChatOpenAI

    # Set up paths and initialize components
    syllabus_path = Path("/path/to/syllabus.docx")
    reading_dir = Path("/path/to/lesson/readings")
    project_dir = Path("/path/to/project")
    llm = ChatOpenAI(api_key="your_api_key")

    # Initialize lesson loader and quiz maker
    lesson_loader = LessonLoader(syllabus_path=syllabus_path, reading_dir=reading_dir, project_dir=project_dir)
    quiz_maker = QuizMaker(
        llm=llm,
        lesson_no=1,
        course_name="Sample Course",
        lesson_loader=lesson_loader,
        output_dir=Path("/path/to/output/dir")
    )

    # Generate and save a quiz
    quiz = quiz_maker.make_a_quiz()
    quiz_maker.save_quiz(quiz)
    quiz_maker.save_quiz_to_ppt(quiz)

"""


# base libraries
import json
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

import pandas as pd
# embedding check for similarity against true questions
import torch
# env setup
from dotenv import load_dotenv
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.prompts import PromptTemplate
# ppt setup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pyprojroot.here import here
# llm chain setup
from sentence_transformers import SentenceTransformer, util

from class_factory.quiz_maker.quiz_prompts import quiz_prompt
from class_factory.quiz_maker.quiz_to_app import quiz_app
from class_factory.quiz_maker.quiz_viz import (generate_dashboard,
                                               generate_html_report)
from class_factory.utils.base_model import BaseModel
from class_factory.utils.llm_validator import Validator
from class_factory.utils.load_documents import LessonLoader
from class_factory.utils.response_parsers import Quiz, ValidatorResponse
from class_factory.utils.tools import retry_on_json_decode_error

load_dotenv()

OPENAI_KEY = os.getenv('openai_key')
OPENAI_ORG = os.getenv('openai_org')

# Path definitions
user_home = Path.home()
wd = here()

syllabus_path = user_home / os.getenv('syllabus_path')
readingDir = user_home / os.getenv("readingsDir")

outputDir = wd / "data/quizzes/"

device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
# %%


# QuizMaker class definition
class QuizMaker(BaseModel):
    """
    A class to generate and manage quizzes based on lesson readings and objectives using a language model (LLM).

    QuizMaker generates quiz questions from lesson content, checks for similarity with prior quizzes to avoid redundancy,
    and validates question format. Quizzes can be saved as Excel or PowerPoint files, launched interactively, and analyzed for performance.

    Attributes:
        llm: The language model instance for quiz generation.
        syllabus_path (Path): Path to the syllabus file.
        reading_dir (Path): Directory containing lesson readings.
        output_dir (Path): Directory for saving quiz files.
        prior_quiz_path (Path): Directory with prior quizzes for similarity checks.
        lesson_range (range): Range of lessons for quiz generation.
        course_name (str): Name of the course for context in question generation.
        device: Device for embeddings (CPU or GPU).
        rejected_questions (List[Dict]): List of questions flagged as similar to prior quizzes.

    Methods:
        make_a_quiz(difficulty_level: int = 5, flag_threshold: float = 0.7) -> List[Dict]:
            Generates quiz questions with similarity checks to avoid redundancy.

        save_quiz(quiz: List[Dict]) -> None:
            Saves quiz questions to an Excel file.

        save_quiz_to_ppt(quiz: List[Dict] = None, excel_file: Path = None, template_path: Path = None) -> None:
            Saves quiz questions to a PowerPoint file, optionally with a template.

        launch_interactive_quiz(quiz_data, sample_size: int = 5, seed: int = 42, save_results: bool = False, output_dir: Path = None, qr_name: str = None) -> None:
            Launches an interactive quiz using Gradio.

        assess_quiz_results(quiz_data: pd.DataFrame = None, results_dir: Path = None, output_dir: Path = None) -> pd.DataFrame:
            Analyzes quiz results and generates summary statistics and visualizations.

    Internal Methods:
        _validate_llm_response(quiz_questions: Dict[str, Any], objectives: str, readings: str, prior_quiz_questions: List[str], difficulty_level: int, additional_guidance: str) -> Dict[str, Any]:
            Validates generated quiz questions for relevance and format.

        _validate_questions(questions: List[Dict]) -> List[Dict]:
            Checks for formatting errors and corrects them.

        _build_quiz_chain() -> Any:
            Builds the LLM chain for quiz generation.

        _load_and_merge_prior_quizzes() -> Tuple[List[str], pd.DataFrame]:
            Loads and merges questions from prior quizzes for similarity checking.

        _check_question_similarity(generated_questions: List[str], threshold: float = 0.6) -> List[Dict]:
            Checks for question similarity against prior quizzes.

        _separate_flagged_questions(questions: List[Dict], flagged_questions: List[Dict]) -> Tuple[List[Dict], List[Dict]]:
            Separates flagged questions based on similarity results.
    """

    def __init__(self, llm, lesson_no: int, course_name: str, lesson_loader: LessonLoader,
                 output_dir: Union[Path, str] = None, prior_quiz_path: Union[Path, str] = None,
                 lesson_range: range = range(1, 5), quiz_prompt_for_llm: str = None, device=None,
                 lesson_objectives: dict = None, verbose=False):
        """
        Initialize QuizMaker with lesson paths and configuration options.

        Args:
            llm: The language model for quiz generation.
            lesson_no (int): Current lesson number.
            course_name (str): Course name for context.
            lesson_loader (LessonLoader): Utility for loading lesson objectives and readings.
            output_dir (Union[Path, str], optional): Directory to save quizzes. Defaults to None.
            prior_quiz_path (Union[Path, str], optional): Directory for prior quizzes to check for question similarity. Defaults to None.
            lesson_range (range): Range of lessons for quiz generation. Defaults to range(1, 5).
            quiz_prompt_for_llm (str): LLM prompt template for generating questions. If not provided, reverts to module default prompt.
            device (optional): Device for embeddings (CPU or GPU). Defaults to CPU if not specified.
            lesson_objectives (optional, dict): user-provided lesson objectives if syllabus not available.
            verbose (bool): Enables verbose logging. Defaults to False.
        """
        # Initialize BaseModel to set up lesson_loader, paths, and logging
        super().__init__(lesson_no=lesson_no, course_name=course_name,
                         lesson_loader=lesson_loader, output_dir=output_dir, verbose=verbose)

        self.llm = llm
        self.prior_quiz_path = self.lesson_loader._validate_dir_path(prior_quiz_path, "prior quiz path")
        self.lesson_range = lesson_range
        self.user_objectives = self.set_user_objectives(lesson_objectives, self.lesson_range) if lesson_objectives else {}

        # Initialize validator, parser, and similarity model
        self.quiz_parser = JsonOutputParser(pydantic_object=Quiz)
        self.val_parser = JsonOutputParser(pydantic_object=ValidatorResponse)
        self.quiz_prompt = quiz_prompt_for_llm if quiz_prompt_for_llm else quiz_prompt
        self.validator = Validator(llm=self.llm, parser=self.val_parser, log_level=self.logger.level)

        # Set device for similarity checking (default to GPU if available)
        self.device = torch.device("cuda" if torch.cuda.is_available() else "cpu") if device is None else device

        # Load prior quiz questions for similarity checking
        self.prior_quiz_questions, self.prior_quizzes = self._load_and_merge_prior_quizzes()
        self.readings = self._load_readings(self.lesson_range)

        # Similarity measures for generated questions
        self.model = SentenceTransformer('all-MiniLM-L6-v2').to(self.device)
        self.rejected_questions = []
        self.generated_questions = []

    @retry_on_json_decode_error()
    def make_a_quiz(self, difficulty_level: int = 5, flag_threshold: float = 0.7) -> List[Dict]:
        """
        Generate quiz questions based on lesson readings and objectives, checking for similarity with prior quizzes.

        Args:
            difficulty_level (int): Difficulty of generated questions, scale 1-10. Defaults to 5.
            flag_threshold (float): Similarity threshold for rejecting duplicate questions. Defaults to 0.7.

        Returns:
            List[Dict[str, Any]]: Generated quiz questions, with duplicates removed. Each dict contains:
                - question (str): The question text
                - type (str): Question type (e.g. "multiple_choice")
                - A) (str): First answer choice
                - B) (str): Second answer choice
                - C) (str): Third answer choice
                - D) (str): Fourth answer choice
                - correct_answer (str): Letter of correct answer
        """
        # Leverage `_load_readings` and `extract_lesson_objectives`
        objectives_dict = {str(lesson): self._get_lesson_objectives(lesson) for lesson in self.lesson_range}
        questions_not_to_use = list(set(self.prior_quiz_questions + [q['question'] for q in self.rejected_questions]))

        chain = self._build_quiz_chain()
        responselist = []

        for lesson_no, readings in self.readings.items():
            self.logger.info(f"\nProcessing Lesson {lesson_no}\n")
            objectives_text = objectives_dict.get(str(lesson_no), "No objectives available")

            # Generate questions using the LLM, building in automatic validation
            additional_guidance = ""
            retries, MAX_RETRIES = 0, 3

            for reading in readings:
                valid = False
                while not valid and retries < MAX_RETRIES:
                    response = chain.invoke({
                        "course_name": self.course_name,
                        "objectives": objectives_text,
                        "information": reading,
                        'prior_quiz_questions': questions_not_to_use,
                        'difficulty_level': difficulty_level,
                        "additional_guidance": additional_guidance
                    })
                    self.logger.info(f"Response from LLM: {response}")
                    # if string, remove the code block indicators (```json and ``` at the end)
                    quiz_questions = json.loads(response.replace('```json\n', '').replace('\n```', '')) if isinstance(response, str) else response

                    val_response = self._validate_llm_response(quiz_questions=quiz_questions,
                                                               objectives=objectives_text,
                                                               reading=reading,
                                                               prior_quiz_questions=questions_not_to_use,
                                                               difficulty_level=difficulty_level,
                                                               additional_guidance=additional_guidance)

                    self.validator.logger.debug(f"Validation output: {val_response}")
                    if int(val_response['status']) == 1:
                        valid = True
                    else:
                        retries += 1
                        additional_guidance = val_response.get("additional_guidance", "")
                        self.validator.logger.warning(f"Lesson {lesson_no}: Response validation failed on attempt {retries}. "
                                                      f"Guidance for improvement: {additional_guidance}")

                # Handle validation failure after max retries
                if not valid:
                    raise ValueError("Validation failed after max retries. Ensure correct prompt and input data. Consider trying a different LLM.")

                responselist.extend(self._parse_llm_questions(quiz_questions))
                self.logger.debug(f"responselist: {responselist}")
        responselist = self._validate_question_format(responselist)

        # Check for similarities
        if self.prior_quiz_questions:
            # Extract just the question text for similarity checking
            generated_question_texts = [q['question'] for q in responselist]
            flagged_questions = self._check_question_similarity(generated_question_texts, threshold=flag_threshold)
            self.logger.info(f"Flagged questions: {flagged_questions}")
            # Separate flagged and scrubbed questions
            scrubbed_questions, flagged_list = self._separate_flagged_questions(responselist, flagged_questions)
            self.rejected_questions.extend(flagged_list)

            self.generated_questions.extend(scrubbed_questions)
            return scrubbed_questions
        else:
            self.generated_questions.extend(responselist)
            return responselist

    def _parse_llm_questions(self, quiz_questions: Union[Dict, List]) -> List[Dict]:
        """
        Process LLM-generated questions for consistent structure and formatting.

        Args:
            quiz_questions (Union[Dict, List]): Questions generated by the LLM, structured as a dictionary by type or a list.

        Returns:
            List[Dict]: Processed list of quiz questions, with each question labeled by type.
        """
        processed_questions = []
        if isinstance(quiz_questions, dict):
            for question_type, questions in quiz_questions.items():
                for question in questions:
                    question['type'] = question_type
                    processed_questions.append(question)
        elif isinstance(quiz_questions, list):
            processed_questions.extend(quiz_questions)
        return processed_questions

    def _validate_llm_response(self, quiz_questions: Dict[str, Any], objectives: str, reading: str,
                               prior_quiz_questions: List[str], difficulty_level: int, additional_guidance: str) -> Dict[str, Any]:
        """
        Validate generated quiz questions by checking relevance and formatting.

        Args:
            quiz_questions (Dict[str, Any]): LLM-generated quiz questions.
            objectives (str): Current lesson objectives.
            readings (str): Lesson readings text for context.
            prior_quiz_questions (List[str]): Previous quiz questions for similarity checks.
            difficulty_level (int): Desired difficulty level (1-10).
            additional_guidance (str): Additional validation instructions.

        Returns:
            Dict[str, Any]: Validation response including evaluation score, status, and suggested improvements.
        """
        # Validate quiz quality and accuracy
        response_str = json.dumps(quiz_questions).replace("{", "{{").replace("}", "}}")
        validation_prompt = self.quiz_prompt.format(course_name=self.course_name,
                                                    objectives=objectives,
                                                    information=reading,
                                                    prior_quiz_questions=prior_quiz_questions,
                                                    difficulty_level=difficulty_level,
                                                    additional_guidance=additional_guidance
                                                    ).replace("{", "{{").replace("}", "}}")

        val_response = self.validator.validate(task_description=validation_prompt,
                                               generated_response=response_str,
                                               min_eval_score=8,
                                               specific_guidance="Evaluate the quality of both the question and the answer choices for each question.")

        return val_response

    def _validate_question_format(self, questions: List[Dict]) -> List[Dict]:
        """
        Validate and correct formatting issues in quiz questions.

        Args:
            questions (List[Dict]): List of generated quiz questions.

        Returns:
            List[Dict]: Corrected list of valid questions.
        """
        validated_questions = []
        rejected_questions = []
        seen_questions = set()  # Track unique question text

        for question_num, question in enumerate(questions):
            try:
                # Track the question text to identify duplicates
                question_text = question.get('question', '').strip()
                if question_text in seen_questions:
                    self.logger.warning(f"Duplicate question detected at {question_num}: {question_text}. Skipping.")
                    continue  # Skip duplicates
                seen_questions.add(question_text)

                # Validate and correct the correct_answer field
                correct_answer = question.get('correct_answer', '').strip()
                if correct_answer not in ['A', 'B', 'C', 'D']:
                    self.logger.warning(f"Incorrect formatting with question {question_num}: {correct_answer}. Attempting to correct.")
                    matched = False
                    for option_letter in ['A', 'B', 'C', 'D']:
                        option_text = question.get(f'{option_letter})', '').strip()
                        if option_text and correct_answer and correct_answer.lower().replace(' ', '') == option_text.lower().replace(' ', ''):
                            question['correct_answer'] = option_letter
                            self.logger.info(f"Corrected formatting error for question {question_num}.")
                            matched = True
                            break
                    if not matched:
                        self.logger.warning(
                            f"Could not match correct_answer '{correct_answer}' to any option in question {question_num}. Adding to rejected questions."
                        )
                        rejected_questions.append(question)
                        continue  # Skip adding this question to the validated list

                # Append validated question
                validated_questions.append(question)

            except Exception as e:
                self.logger.error(f"Error validating question {question_num}: {e}. Adding to rejected questions.")
                rejected_questions.append(question)

        # Optionally log or store rejected questions for debugging
        if rejected_questions:
            self.rejected_questions.extend(rejected_questions)  # Assuming self.rejected_questions exists
            self.logger.info(f"Rejected {len(rejected_questions)} questions due to validation issues.")

        return validated_questions

    def _build_quiz_chain(self):
        """
        Construct the LLM quiz generation chain using prompt, LLM, and output parser.

        Returns:
            Chain: Quiz generation chain combining prompt template, LLM, and parser.
        """
        # Create the chain with the prompt, LLM, and parser
        chain = self.quiz_prompt | self.llm | self.quiz_parser

        return chain

    def _load_and_merge_prior_quizzes(self) -> Tuple[List[str], pd.DataFrame]:
        """
        Load and merge previous quiz questions from Excel files for similarity checking.

        Returns:
            Tuple[List[str], pd.DataFrame]: List of prior quiz questions and a DataFrame of prior quizzes.
        """
        all_quiz_data = []
        quiz_df = pd.DataFrame()

        # Get all Excel files in the prior_quiz_dir
        for file in self.prior_quiz_path.glob('*.xlsx'):
            df = pd.read_excel(file)
            quiz_df = pd.concat([quiz_df, df], axis=0, ignore_index=True)
            if 'question' in df.columns:
                all_quiz_data.append(df['question'].tolist())

        # Flatten the list of lists into a single list
        merged_questions = [item for sublist in all_quiz_data for item in sublist]
        quiz_df = quiz_df.reset_index(drop=True)

        return merged_questions, quiz_df

    # Function to check similarity between generated questions and main quiz questions
    def _check_question_similarity(self, generated_questions: List[str], threshold: float = 0.6) -> List[Dict]:
        """
        Compare generated questions to prior quiz questions to flag similar questions.

        Args:
            generated_questions (List[str]): List of new quiz questions.
            threshold (float): Similarity threshold for flagging. Defaults to 0.6.

        Returns:
            List[Dict]: List of flagged questions with similarity scores and indices.
        """
        flagged_questions = []

        # Get embeddings for the newly generated questions and prior quiz questions
        generated_embeddings = self.model.encode(generated_questions, convert_to_tensor=True, device=self.device)
        prior_quiz_embeddings = self.model.encode(self.prior_quiz_questions, convert_to_tensor=True, device=self.device)

        # Compare each generated question to prior quiz questions
        for i, gen_question in enumerate(generated_questions):
            cosine_scores = util.pytorch_cos_sim(generated_embeddings[i], prior_quiz_embeddings)
            max_score = float(cosine_scores.max())

            if max_score > threshold:
                flagged_questions.append({'index': i, 'question': gen_question, 'similarity': max_score})

        return flagged_questions

    def _separate_flagged_questions(self, questions: List[Dict], flagged_questions: List[Dict]) -> Tuple[List[Dict], List[Dict]]:
        """
        Separates flagged questions from those that pass similarity checks.

        Args:
            questions (List[Dict]): List of generated quiz questions.
            flagged_questions (List[Dict]): List of questions flagged as similar.

        Returns:
            Tuple[List[Dict], List[Dict]]: List of valid questions and list of flagged questions.
        """
        scrubbed_questions = []
        flagged_list = []
        unique_questions = set()

        flagged_indices = {f['index'] for f in flagged_questions}  # Get set of flagged indices

        for i, question in enumerate(questions):
            if i in flagged_indices:
                flagged_list.append(question)
            else:
                if question['question'].strip() in unique_questions:
                    continue
                unique_questions.add(question['question'].strip())
                scrubbed_questions.append(question)

        return scrubbed_questions, flagged_list

    def save_quiz(self, quiz: List[Dict]) -> None:
        """
        Save quiz questions to an Excel file.

        Args:
            quiz (List[Dict[str, Any]]): List of quiz questions to save. Each dict should contain:
                - type (str): Question type
                - question (str): Question text
                - A) (str): First answer choice
                - B) (str): Second answer choice
                - C) (str): Third answer choice
                - D) (str): Fourth answer choice
                - correct_answer (str): Letter of correct answer

        Returns:
            None
        """
        final_quiz = pd.DataFrame(quiz)

        col_order = ['type', 'question', 'A)', 'B)', 'C)', 'D)', 'correct_answer']
        final_quiz = final_quiz[col_order].reset_index(drop=True)
        final_quiz.to_excel(self.output_dir / f"l{min(self.lesson_range)}_{max(self.lesson_range)}_quiz.xlsx", index=False)

    def save_quiz_to_ppt(self, quiz: List[Dict] = None, excel_file: Union[Path, str] = None,
                         template_path: Union[Path, str] = None, filename: str = None) -> None:
        """
        Save quiz questions to a PowerPoint presentation, with options to use a template.

        Args:
            quiz (List[Dict], optional): List of quiz questions in dictionary format.
            excel_file (Path, optional): Path to an Excel file containing quiz questions. If provided, it overrides the quiz argument.
            template_path (Path, optional): Path to a PowerPoint template to apply to the generated slides.

        Raises:
            ValueError: If neither quiz nor excel_file is provided.

        Creates a PowerPoint presentation with each question on a slide, followed by the answer slide.
        """
        # Load from Excel if an Excel file path is provided
        if excel_file:
            excel_file = Path(excel_file)
            df = pd.read_excel(excel_file).fillna('')
        elif quiz:
            # Convert the list of dicts into a DataFrame
            df = pd.DataFrame(quiz)
        else:
            raise ValueError("Either 'quiz' or 'excel_file' must be provided.")

        use_template = template_path and template_path.exists()

        if use_template:
            template_path = Path(template_path)
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()  # Use default template if none is provided

        # Function to add a slide with a title and content
        def _add_slide(prs, title, content, answer=None, is_answer=False, bg_color=(255, 255, 255)):
            slide_layout = prs.slide_layouts[1]  # Layout with title and content
            slide = prs.slides.add_slide(slide_layout)

            if not use_template:
                # Set slide background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])

            # Set title text and font size
            title_placeholder = slide.shapes.title
            title_placeholder.text = title
            if not use_template:
                title_text_frame = title_placeholder.text_frame
                title_text_frame.paragraphs[0].font.size = Pt(24)  # Set title font size

            # Set content text and font size
            text_placeholder = slide.shapes.placeholders[1]
            text_placeholder.text = content
            if not use_template:
                content_text_frame = text_placeholder.text_frame
                content_text_frame.paragraphs[0].font.size = Pt(32)  # Set content font size

            if is_answer and answer:
                # For answers, add the answer text at the end and adjust the font size
                text_placeholder.text += f"\n\nAnswer: {answer}"
                if not use_template:
                    content_text_frame.paragraphs[0].font.size = Pt(32)  # Adjust answer font size

        # Loop through each question and add to the presentation
        for i, row in df.iterrows():
            # Get question text
            question_text = row['question']

            # Prepare choices, using .get() to avoid KeyError
            choices = f"A) {row['A)']}\nB) {row['B)']}"
            if row.get('C)', '').strip():
                choices += f"\nC) {row['C)']}"
            if row.get('D)', '').strip():
                choices += f"\nD) {row['D)']}"

            # Add a slide for the question
            _add_slide(prs, f"Question {i + 1}", question_text + "\n\n" + choices, bg_color=(255, 255, 255))

            # Add a slide for the answer
            correct_answer = row['correct_answer']
            if correct_answer in ['A', 'B', 'C', 'D']:
                answer_text = row[f'{correct_answer})']
                _add_slide(prs, f"Answer to Question {i + 1}", question_text,
                           answer=f"{correct_answer}: {answer_text}", is_answer=True, bg_color=(255, 255, 255))
            else:
                _add_slide(prs, f"Answer to Question {i + 1}", question_text,
                           answer=f"{correct_answer}", is_answer=True, bg_color=(255, 255, 255))

        # Function to remove the slide at index 0 (first slide). Only applicable if using a template
        def _remove_slide_by_index(pres, index):
            """Remove the slide at the specified index. Using a template leaves the first slide blank."""
            slide_id = pres.slides._sldIdLst[index]
            pres.part.drop_rel(slide_id.rId)
            pres.slides._sldIdLst.remove(slide_id)

        # Remove the first slide if using a template
        if use_template:
            _remove_slide_by_index(prs, 0)

        # Save the PowerPoint presentation
        if filename:
            ppt_path = self.output_dir / f"{filename}.pptx"
        elif excel_file:
            ppt_path = excel_file.with_suffix(".pptx")
        else:
            ppt_path = self.output_dir / f"quiz_presentation_{min(self.lesson_range)}_{max(self.lesson_range)}.pptx"
        prs.save(ppt_path)
        print(f"Presentation saved at {ppt_path}")

    def launch_interactive_quiz(
        self,
        quiz_data: Union[pd.DataFrame, Path, str, List[Dict[str, Any]]] = None,
        sample_size: int = 5,
        seed: int = 42,
        save_results: bool = False,
        output_dir: Optional[Path] = None,
        qr_name: Optional[str] = None
    ) -> None:
        """
        Launch an interactive quiz using Gradio, sampling questions from provided data or generating new data if none is provided.

        Args:
            quiz_data (Union[pd.DataFrame, Path, str, List[Dict[str, Any]]], optional): Quiz questions as a DataFrame,
                Excel path, or list of dictionaries. If None, generates new questions.
            sample_size (int, optional): Number of questions to sample. Defaults to 5.
            seed (int, optional): Random seed for consistent sampling. Defaults to 42.
            save_results (bool, optional): Whether to save quiz results. Defaults to False.
            output_dir (Path, optional): Directory to save quiz results. Defaults to the class’s output directory.
            qr_name (str, optional): Name of the QR code image file for accessing the quiz.

        Raises:
            ValueError: If quiz_data is provided but is not a valid type.
        """

        # If quiz_data is a List of Dicts, convert it to a DataFrame
        if isinstance(quiz_data, list) and isinstance(quiz_data[0], dict):  # Check if it's a List[Dict]
            quiz = pd.DataFrame(quiz_data)

        # If quiz_data is a path or string, read the file
        elif isinstance(quiz_data, (Path, str)):
            quiz = pd.read_excel(quiz_data)

        # If quiz_data is None, generate the quiz dynamically
        elif quiz_data is None:
            self.logger.warning("No quiz identified. Generating new quiz from lesson documents")
            quiz_generated = self.make_a_quiz()
            quiz = pd.DataFrame(quiz_generated)

        # If quiz_data is already a DataFrame, use it directly
        elif isinstance(quiz_data, pd.DataFrame):
            quiz = quiz_data

        else:
            raise ValueError("Invalid type for quiz_data. It must be either a DataFrame, path to an Excel file, or a list of dictionaries.")

        quiz_sampled = quiz.sample(sample_size, random_state=seed)

        if not output_dir:
            output_dir = self.output_dir

        quiz_app(quiz_sampled, save_results=save_results,
                 output_dir=output_dir, qr_name=qr_name)

    def assess_quiz_results(
        self,
        quiz_data: Optional[pd.DataFrame] = None,
        results_dir: Optional[Union[Path, str]] = None,
        output_dir: Optional[Union[Path, str]] = None
    ) -> pd.DataFrame:
        """
        Analyze quiz results, generate summary statistics, and visualize responses.

        Args:
            quiz_data (pd.DataFrame, optional): DataFrame of quiz results. If None, loads results from CSV files in results_dir.
            results_dir (Path, optional): Directory containing CSV files of quiz results.
            output_dir (Path, optional): Directory for saving summary statistics and plots. Defaults to output_dir/'quiz_analysis'.

        Returns:
            pd.DataFrame: DataFrame with summary statistics, including:
                - question (str): Question text
                - Total Responses (int): Number of unique users who answered
                - Correct Responses (int): Number of correct answers
                - Incorrect Responses (int): Number of incorrect answers
                - Percent Correct (float): Percentage of correct answers
                - Modal Answer (str): Most common answer given

        Raises:
            AssertionError: If quiz_data is provided but is not a pandas DataFrame.
        """
        # Determine the output directory
        if results_dir:
            results_dir = Path(results_dir)
            output_dir = results_dir.parent / 'quiz_analysis'
        else:
            output_dir = Path(output_dir) if output_dir else self.output_dir / 'quiz_analysis'

        # If quiz_data is provided as a DataFrame, use it directly
        if quiz_data is not None:
            assert isinstance(quiz_data, pd.DataFrame), "if passing an object for assessment, it needs to be a pd.DataFrame"
            df = quiz_data
        else:
            # Load from CSV if no DataFrame is provided
            results_dir = Path(results_dir) if results_dir else self.output_dir / 'quiz_results'

            csv_files = list(results_dir.glob('*.csv'))
            if not csv_files:
                self.logger.warning(f"No CSV files found in directory: {results_dir}")
                return pd.DataFrame()

            df_list = [pd.read_csv(f) for f in csv_files]
            df = pd.concat(df_list, ignore_index=True)

        output_dir.mkdir(parents=True, exist_ok=True)

        # Sort by timestamp to ensure records are in chronological order
        df = df.sort_values(by='timestamp')

        # Remove duplicates keeping last attempt for each user per question
        df = df.drop_duplicates(subset=['user_id', 'question'], keep='last').reset_index(drop=True)
        # Ensure 'is_correct' is boolean
        df['is_correct'] = df['is_correct'].astype(bool)

        # Calculate summary statistics
        summary = df.groupby('question').apply(lambda x: pd.Series({
            'Total Responses': x['user_id'].nunique(),
            'Correct Responses': x[x['is_correct'] == True]['user_id'].nunique(),
            'Incorrect Responses': x[x['is_correct'] == False]['user_id'].nunique(),
            'Percent Correct': x['is_correct'].mean() * 100,
            'Modal Answer': x['user_answer'].mode()[0] if not x['user_answer'].mode().empty else None
        }), include_groups=False).reset_index()

        # Save summary statistics to a CSV file
        summary_output_path = output_dir / 'summary_statistics.csv'
        summary.to_csv(summary_output_path, index=False)

        # Generate the HTML report
        generate_html_report(df, summary, output_dir)
        # Generate the dashboard
        generate_dashboard(df, summary)

        return summary


# %%
if __name__ == "__main__":

    from langchain_anthropic import ChatAnthropic
    from langchain_community.llms import Ollama
    from langchain_openai import ChatOpenAI
    from pyprojroot.here import here

    from class_factory.utils.tools import reset_loggers
    ANTHROPIC_API_KEY = os.getenv("anthropic_api_key")

    reset_loggers()

    user_home = Path.home()
    wd = here()

    # llm = ChatOpenAI(
    #     model="gpt-4o-mini",
    #     temperature=0.8,
    #     max_tokens=None,
    #     timeout=None,
    #     max_retries=2,
    #     api_key=OPENAI_KEY,
    #     organization=OPENAI_ORG,
    # )
    llm = ChatAnthropic(
        model="claude-3-5-haiku-latest",
        temperature=0.5,
        max_retries=2,
        api_key=ANTHROPIC_API_KEY
    )
    # llm = Ollama(
    #     model="mistral", # 4k context window
    #     #model="yarn-mistral", # for 64k context window, also could use yarn-mistral:7b-128k if system memory allows
    #     #model="llama3.1",
    #     temperature=0.0,
    #     format="json"
    # )

    lesson_no = 20

    # Path definitions
    readingDir = user_home / os.getenv('readingsDir')
    slideDir = user_home / os.getenv('slideDir')
    syllabus_path = user_home / os.getenv('syllabus_path')

    loader = LessonLoader(syllabus_path=syllabus_path,
                          reading_dir=readingDir,
                          slide_dir=slideDir)

    maker = QuizMaker(llm=llm,
                      lesson_loader=loader,
                      output_dir=wd/"ClassFactoryOutput/QuizMaker",
                      quiz_prompt_for_llm=quiz_prompt,
                      lesson_no=20,
                      lesson_range=range(19, 21),
                      course_name="American Government",
                      prior_quiz_path=outputDir,
                      verbose=True)

    # quiz_name = wd / f"ClassFactoryOutput/QuizMaker/L24/l22_24_quiz.xlsx"
    # quiz_path = wd / f"ClassFactoryOutput/QuizMaker/"

    quiz = maker.make_a_quiz(difficulty_level=0.9)
    # maker.save_quiz_to_ppt(quiz)

    # maker.save_quiz(quiz)
    # maker.launch_interactive_quiz(quiz_name, sample_size=5, save_results=True,
    #                               output_dir=quiz_path, qr_name="test_quiz")
    # maker.assess_quiz_results()
