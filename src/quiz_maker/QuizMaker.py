"""
**QuizMaker**

This module defines the `QuizMaker` class, which generates quiz questions based on lesson readings and objectives using a language model (LLM).
The generated quizzes can be saved in Excel or PowerPoint formats, and the class also supports launching interactive quizzes with a web interface,
analyzing quiz results, and avoiding duplication of questions from prior quizzes.

The class has the following functionalities:

- **Quiz Generation**: Automatically generates quiz questions based on specified lesson objectives and readings.

- **Similarity Checking**: Ensures generated questions do not overlap significantly with prior quiz questions, using sentence embedding models for similarity checks.

- **Question Validation**: Validates and corrects the format of generated quiz questions to ensure that answers are in the proper format (e.g., 'A', 'B', 'C', 'D').

- **Saving**: Quizzes can be saved as Excel files or converted into PowerPoint presentations, with support for PowerPoint templates.

- **Interactive Quiz Launch**: Integrates with Gradio to launch interactive quizzes that users can take in real time, with results saved and analyzed.

- **Results Assessment**: Analyzes quiz results from CSV files, calculating summary statistics and generating visual reports and dashboards.

Dependencies:

- `langchain_core`: For prompt and LLM integration.
- `sentence_transformers`: For sentence embeddings used in similarity checks.
- `pptx`: For generating PowerPoint presentations.
- `pandas`: For data handling and analysis.
- `torch`: For managing device usage (CPU or GPU) and embeddings.
- `gradio`: For interactive quiz functionality.
- Custom utility modules for document loading, response parsing, logging, and retry decorators.

Usage:

1. **Quiz Generation**:

   Instantiate `QuizMaker` with the necessary paths and LLM, then call `make_a_quiz()` to generate the quiz.

2. **Saving Quizzes**:

   After generating quiz questions, use the `save_quiz()` method to save the quiz as an Excel file or `save_quiz_to_ppt()`
   to save it as a PowerPoint presentation. You can optionally use a PowerPoint template for custom slide design.

3. **Interactive Quiz Launch**:

   Use `launch_interactive_quiz()` to launch an interactive quiz interface via Gradio. This allows users to participate in quizzes,
   with options to save results and display a QR code for access.

4. **Similarity Checking**:

   During quiz generation, the `make_a_quiz()` method checks for similarity between generated questions and prior quizzes
   to avoid question leakage. The similarity is checked using sentence embeddings, and flagged questions are removed.

5. **Question Validation**:

   The `validate_questions()` method ensures that generated quiz questions have the correct format and that the correct answer
   is properly aligned with the answer choices.

6. **Results Assessment**:

   After conducting a quiz, use the `assess_quiz_results()` method to load CSV files containing user responses and
   generate summary statistics. The method also provides visual analysis, including HTML reports and interactive dashboards.

"""


# base libraries
import json
import logging
import os
from pathlib import Path
from typing import Dict, List, Tuple, Union

import pandas as pd
# embedding check for similarity against true questions
import torch
# env setup
from dotenv import load_dotenv
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.prompts import PromptTemplate
# ppt setup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pyprojroot.here import here
# llm chain setup
from sentence_transformers import SentenceTransformer, util

from src.quiz_maker.quiz_prompts import quiz_prompt
from src.quiz_maker.quiz_to_app import quiz_app
from src.quiz_maker.quiz_viz import generate_dashboard, generate_html_report
# self-defined utils
from src.utils.load_documents import extract_lesson_objectives, load_lessons
from src.utils.response_parsers import Quiz
from src.utils.tools import (logger_setup, reset_loggers,
                             retry_on_json_decode_error)

reset_loggers()
load_dotenv()

OPENAI_KEY = os.getenv('openai_key')
OPENAI_ORG = os.getenv('openai_org')

# Path definitions
wd = here()
syllabus_path = Path(os.getenv('syllabus_path'))
readingDir = Path(os.getenv("readingsDir"))

outputDir = wd / "data/quizzes/"

device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
# %%


# QuizMaker class definition
class QuizMaker:
    """
    A class to generate quizzes based on lesson readings and objectives using a language model (LLM).
    This class also allows for easy distribution via QR code and analysis of quiz results.

    Generated quizzes can be saved to excel or powerpoint; and the class supports saving to a specific provided ppt template
    The quiz generation process has the option of uploading prior (or upcoming) quizzes, to ensure real quiz questions aren't leaked to practice quizzes

    Attributes:
        llm: The language model instance for generating quiz questions.
        syllabus_path (Path): Path to the syllabus file.
        reading_dir (Path): Directory where lesson readings are stored.
        output_dir (Path): Directory where the generated quiz will be saved.
        prior_quiz_path (Path): Path to the previous quiz for similarity checking.
        lesson_range (range): The range of lessons for which to generate quizzes.
        course_name (str): Name of the course for quiz generation context.
        device: The device for sentence embeddings (CPU or GPU).
        rejected_questions (List[Dict]): List of questions rejected due to similarity.
    """

    def __init__(self, llm, syllabus_path: Union[Path, str], reading_dir: Union[Path, str],
                 output_dir: Union[Path, str], prior_quiz_path: Union[Path, str],
                 lesson_range: range, quiz_prompt: str = quiz_prompt, device=None,
                 course_name: str = 'Political Science', verbose=False):
        """
        Initialize QuizMaker with the necessary paths, LLM, and other configurations.

        Args:
            llm: The language model instance for generating quiz questions.
            syllabus_path (Path): Path to the syllabus file.
            reading_dir (Path): Directory where lesson readings are stored.
            output_dir (Path): Directory where the generated quiz will be saved.
            prior_quiz_path (Path): Path to the previous quiz for similarity checking.
            lesson_range (range): Range of lessons for which to generate quizzes.
            quiz_prompt (str, optional): The LLM prompt template for generating quiz questions. Defaults to the module-level `quiz_prompt`.
            device (Optional[str], optional): The device for sentence embeddings (CPU or GPU). Defaults to GPU if available.
            course_name (str, optional): The name of the course for quiz generation context. Defaults to 'Political Science'.
            verbose (bool, optional): Whether to output verbose logs. Defaults to False.
        """
        self.llm = llm
        self.syllabus_path = syllabus_path
        self.reading_dir = Path(reading_dir)
        self.output_dir = Path(output_dir)
        self.root_dir = self.output_dir.parent
        self.prior_quiz_path = prior_quiz_path
        self.lesson_range = lesson_range
        self.course_name = course_name

        # setup logging
        log_level = logging.INFO if verbose else logging.WARNING
        self.logger = logger_setup(log_level=log_level)

        # Define the LLM prompt template
        self.quiz_parser = JsonOutputParser(pydantic_object=Quiz)
        self.quiz_prompt = quiz_prompt

        # Set device for similarity checking (default to GPU if available)
        self.device = torch.device("cuda" if torch.cuda.is_available() else "cpu") if device is None else device

        # Load prior quiz questions for similarity checking
        self.prior_quiz_questions, self.prior_quizzes = self.load_and_merge_prior_quizzes()

        # Initialize sentence transformer model for similarity checking
        self.model = SentenceTransformer('all-MiniLM-L6-v2').to(self.device)
        self.rejected_questions = []

    @retry_on_json_decode_error()
    def make_a_quiz(self, flag_threshold: float = 0.7) -> List[Dict]:
        """
        Generate quiz questions based on lesson readings and objectives, for each lesson in the lesson range.

        Args:
            flag_threshold (float): The similarity threshold beyond which a generated quiz question will be rejected. Defaults to 0.7.

        Returns:
            List[Dict]: A list of generated quiz questions, scrubbed for similarity.
        """
        all_readings = []
        objectives = []
        responselist = []

        for lesson_no in self.lesson_range:
            input_dir = self.reading_dir / f'L{lesson_no}/'

            # load_documents to process all readings for the current lesson
            readings = load_lessons(input_dir, lesson_no, recursive=False)
            all_readings.extend(readings)

            # Extract lesson objectives
            objectives_text = extract_lesson_objectives(self.syllabus_path, lesson_no, only_current=True)
            objectives.append(objectives_text)

            # add rejected questions, if any, to our list of things to avoid
            rejected_questions = [q['question'] for q in self.rejected_questions] if self.rejected_questions else []
            questions_not_to_use = list(set(self.prior_quiz_questions + rejected_questions))

            chain = self.build_quiz_chain()
            # Generate questions using the LLM
            response = chain.invoke({
                "course_name": self.course_name,
                "objectives": objectives_text,
                "information": readings,
                'prior_quiz_questions': questions_not_to_use
            })

            # if string, remove the code block indicators (```json and ``` at the end)
            if isinstance(response, str):
                response_cleaned = response.replace('```json\n', '').replace('\n```', '')
                quiz_questions = json.loads(response_cleaned)
            else:
                quiz_questions = response

            if isinstance(quiz_questions, dict):
                # Flatten the questions and add a 'type' key
                for question_type, questions in quiz_questions.items():
                    for question in questions:
                        # Add the question type as a key to each question
                        question['type'] = question_type
                        # Add the updated question to responselist
                        responselist.append(question)
            else:
                responselist.extend(quiz_questions)

        responselist = self.validate_questions(responselist)

        # Extract just the question text for similarity checking
        generated_question_texts = [q['question'] for q in responselist]

        # Check for similarities
        flagged_questions = self.check_question_similarity(generated_question_texts, threshold=flag_threshold)

        # Separate flagged and scrubbed questions
        scrubbed_questions, flagged_list = self.separate_flagged_questions(responselist, flagged_questions)
        self.rejected_questions.extend(flagged_list)

        return scrubbed_questions

    def validate_questions(self, questions: List[Dict]) -> List[Dict]:
        """
        Validate the generated questions and correct formatting issues if found.

        Args:
            questions (List[Dict]): The list of generated quiz questions.

        Returns:
            List[Dict]: The validated and potentially corrected list of quiz questions.
        """
        for question_num, question in enumerate(questions):
            correct_answer = question.get('correct_answer', '').strip()
            if correct_answer not in ['A', 'B', 'C', 'D']:
                self.logger.warning(f"Incorrect formatting with {question_num}: {correct_answer}. Attempting to correct.")
                matched = False
                for option_letter in ['A', 'B', 'C', 'D']:
                    option_text = question.get(f'{option_letter})', '').strip()
                    if option_text and correct_answer and correct_answer.lower().replace(' ', '') == option_text.lower().replace(' ', ''):
                        question['correct_answer'] = option_letter
                        matched = True
                        break
                if not matched:
                    self.logger.warning(
                        f"Could not match correct_answer '{correct_answer}' to any option in question '{question.get('question', '')}'. Recommend manual fix.")
                    question['correct_answer'] = None
        return questions

    def build_quiz_chain(self):
        """
        Build the quiz generation chain by combining the prompt, LLM, and parser.

        Returns:
            The quiz generation chain.
        """
        # Construct the prompt template
        combined_template = PromptTemplate.from_template(self.quiz_prompt)
        # Create the chain with the prompt, LLM, and parser
        chain = combined_template | self.llm | self.quiz_parser

        return chain

    def load_and_merge_prior_quizzes(self) -> list:
        """
        Load and merge all prior quiz questions from Excel files in the prior_quiz_dir.

        Returns:
            Tuple[List[str], pd.DataFrame]: A list of prior quiz questions and a DataFrame of prior quizzes.
        """
        all_quiz_data = []
        quiz_df = pd.DataFrame()

        # Get all Excel files in the prior_quiz_dir
        for file in self.prior_quiz_path.glob('*.xlsx'):
            df = pd.read_excel(file)
            quiz_df = pd.concat([quiz_df, df], axis=0, ignore_index=True)
            if 'question' in df.columns:
                all_quiz_data.append(df['question'].tolist())

        # Flatten the list of lists into a single list
        merged_questions = [item for sublist in all_quiz_data for item in sublist]
        quiz_df = quiz_df.reset_index(drop=True)

        return merged_questions, quiz_df

    # Function to check similarity between generated questions and main quiz questions
    def check_question_similarity(self, generated_questions, threshold=0.6) -> List[Dict]:
        """
        Check similarity between generated quiz questions and prior quiz questions.

        Args:
            generated_questions (List[str]): List of generated quiz questions.
            threshold (float): The similarity threshold for flagging questions. Defaults to 0.6.

        Returns:
            List[Dict]: A list of flagged questions with similarity scores and their indices.
        """
        flagged_questions = []

        # Get embeddings for the newly generated questions and prior quiz questions
        generated_embeddings = self.model.encode(generated_questions, convert_to_tensor=True, device=self.device)
        prior_quiz_embeddings = self.model.encode(self.prior_quiz_questions, convert_to_tensor=True, device=self.device)

        # Compare each generated question to prior quiz questions
        for i, gen_question in enumerate(generated_questions):
            cosine_scores = util.pytorch_cos_sim(generated_embeddings[i], prior_quiz_embeddings)
            max_score = float(cosine_scores.max())

            if max_score > threshold:
                flagged_questions.append({'index': i, 'question': gen_question, 'similarity': max_score})

        return flagged_questions

    def separate_flagged_questions(self, questions, flagged_questions) -> Tuple[List, List]:
        """
        Validate the generated questions and correct formatting issues if found.

        Args:
            questions (List[Dict]): The list of generated quiz questions.

        Returns:
            List[Dict]: The validated and potentially corrected list of quiz questions.
        """
        scrubbed_questions = []
        flagged_list = []

        flagged_indices = {f['index'] for f in flagged_questions}  # Get set of flagged indices

        for i, question in enumerate(questions):
            if i in flagged_indices:
                flagged_list.append(question)
            else:
                scrubbed_questions.append(question)

        return scrubbed_questions, flagged_list

    def save_quiz(self, quiz: List[Dict]) -> None:
        """Save to excel"""
        final_quiz = pd.DataFrame(quiz)

        col_order = ['type', 'question', 'A)', 'B)', 'C)', 'D)', 'correct_answer']
        final_quiz = final_quiz[col_order].reset_index(drop=True)
        final_quiz.to_excel(self.output_dir / f"l{min(self.lesson_range)}_{max(self.lesson_range)}_quiz.xlsx", index=False)

    def save_quiz_to_ppt(self, quiz: List[Dict] = None, excel_file: Path = None, template_path: Path = None) -> None:
        """
        Save quiz questions to a PowerPoint presentation.

        Args:
            quiz (List[Dict], optional): The list of quiz dictionaries. Defaults to None.
            excel_file (Path, optional): The path to an Excel file containing quiz questions. Defaults to None.
        """
        # Load from Excel if an Excel file path is provided
        if excel_file:
            df = pd.read_excel(excel_file)
        elif quiz:
            # Convert the list of dicts into a DataFrame
            df = pd.DataFrame(quiz)
        else:
            raise ValueError("Either 'quiz' or 'excel_file' must be provided.")

        use_template = template_path and template_path.exists()

        if use_template:
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()  # Use default template if none is provided

        # Function to add a slide with a title and content
        def add_slide(prs, title, content, answer=None, is_answer=False, bg_color=(255, 255, 255)):
            slide_layout = prs.slide_layouts[1]  # Layout with title and content
            slide = prs.slides.add_slide(slide_layout)

            if not use_template:
                # Set slide background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])

            # Set title text and font size
            title_placeholder = slide.shapes.title
            title_placeholder.text = title
            if not use_template:
                title_text_frame = title_placeholder.text_frame
                title_text_frame.paragraphs[0].font.size = Pt(24)  # Set title font size

            # Set content text and font size
            text_placeholder = slide.shapes.placeholders[1]
            text_placeholder.text = content
            if not use_template:
                content_text_frame = text_placeholder.text_frame
                content_text_frame.paragraphs[0].font.size = Pt(32)  # Set content font size

            if is_answer and answer:
                # For answers, add the answer text at the end and adjust the font size
                text_placeholder.text += f"\n\nAnswer: {answer}"
                if not use_template:
                    content_text_frame.paragraphs[0].font.size = Pt(32)  # Adjust answer font size

        # Loop through each question and add to the presentation
        for i, row in df.iterrows():
            # Get question text
            question_text = row['question']

            # Prepare choices (exclude C and D if they are blank)
            choices = f"A) {row['A)']}\nB) {row['B)']}"
            if pd.notna(row['C)']) and row['C)'].strip():  # Only add C if it is not blank
                choices += f"\nC) {row['C)']}"
            if pd.notna(row['D)']) and row['D)'].strip():  # Only add D if it is not blank
                choices += f"\nD) {row['D)']}"

            # Add a slide for the question
            add_slide(prs, f"Question {i + 1}", question_text + "\n\n" + choices, bg_color=(255, 255, 255))

            # Add a slide for the answer
            correct_answer = row['correct_answer']
            if correct_answer in ['A', 'B', 'C', 'D']:
                answer_text = row[f'{correct_answer})']
                add_slide(prs, f"Answer to Question {i + 1}", question_text,
                          answer=f"{correct_answer}: {answer_text}", is_answer=True, bg_color=(255, 255, 255))
            else:
                add_slide(prs, f"Answer to Question {i + 1}", question_text,
                          answer=f"{correct_answer}", is_answer=True, bg_color=(255, 255, 255))

        # Function to remove the slide at index 0 (first slide). Only applicable if using a template
        def remove_slide_by_index(pres, index):
            """Remove the slide at the specified index."""
            slide_id = pres.slides._sldIdLst[index]
            pres.part.drop_rel(slide_id.rId)
            pres.slides._sldIdLst.remove(slide_id)

        # Remove the first slide if using a template
        if use_template:
            remove_slide_by_index(prs, 0)

        # Save the PowerPoint presentation
        ppt_path = self.output_dir / f"quiz_presentation_{min(self.lesson_range)}_{max(self.lesson_range)}.pptx"
        if excel_file:
            ppt_path = excel_file.with_suffix(".pptx")
        prs.save(ppt_path)
        print(f"Presentation saved at {ppt_path}")

    def launch_interactive_quiz(self, quiz_data: Union[pd.DataFrame, Path, str, List[Dict]] = None, sample_size: int = 5, seed: int = 42,
                                save_results: bool = False, output_dir: Path = None, qr_name: str = None) -> None:
        """
        Launch the interactive quiz using Gradio, using either preloaded quiz data or
        dynamically generated quiz data from the class.

        Args:
            quiz_data (Union[pd.DataFrame, Path, str, List[Dict]], optional): The quiz data to be used for the interactive quiz.
                                                                             If not provided, it will use generated quiz data.
            sample_size (int, Default=5): The number of questions to sample from the large df of questions.
        """

        # If quiz_data is a List of Dicts, convert it to a DataFrame
        if isinstance(quiz_data, list) and isinstance(quiz_data[0], dict):  # Check if it's a List[Dict]
            quiz = pd.DataFrame(quiz_data)

        # If quiz_data is a path or string, read the file
        elif isinstance(quiz_data, (Path, str)):
            quiz = pd.read_excel(quiz_data)

        # If quiz_data is None, generate the quiz dynamically
        elif quiz_data is None:
            self.logger.warning("No quiz identified. Generating new quiz from lesson documents")
            quiz_generated = self.make_a_quiz()
            quiz = pd.DataFrame(quiz_generated)

        # If quiz_data is already a DataFrame, use it directly
        elif isinstance(quiz_data, pd.DataFrame):
            quiz = quiz_data

        else:
            raise ValueError("Invalid type for quiz_data. It must be either a DataFrame, path to an Excel file, or a list of dictionaries.")

        quiz_sampled = quiz.sample(sample_size, random_state=seed)

        if not output_dir:
            output_dir = self.output_dir

        quiz_app(quiz_sampled, save_results=save_results,
                 output_dir=output_dir, qr_name=qr_name)

    def assess_quiz_results(self,  output_dir: Path = None) -> pd.DataFrame:
        """
        Load quiz results from CSV files, calculate summary statistics, and generate plots.

        Args:
            results_dir (Path): Directory containing CSV files of quiz results.
            output_dir (Path, optional): Directory where summary and plots will be saved.
                                         If not provided, defaults to self.output_dir / 'quiz_analysis'.

        Returns:
            pd.DataFrame: Summary statistics DataFrame.
        """
        if output_dir is None:
            output_dir = self.output_dir / 'quiz_analysis'
        output_dir.mkdir(parents=True, exist_ok=True)

        # Load all CSV files from the results directory
        results_dir = self.output_dir / 'quiz_results'
        csv_files = list(results_dir.glob('*.csv'))
        if not csv_files:
            self.logger.warning(f"No CSV files found in directory: {results_dir}")
            return pd.DataFrame()

        df_list = [pd.read_csv(f) for f in csv_files]
        df = pd.concat(df_list, ignore_index=True)

        # Ensure 'is_correct' is boolean
        df['is_correct'] = df['is_correct'].astype(bool)

        # Calculate summary statistics
        summary = df.groupby('question').apply(lambda x: pd.Series({
            'Total Responses': x['user_id'].nunique(),
            'Correct Responses': x[x['is_correct'] == True]['user_id'].nunique(),
            'Incorrect Responses': x[x['is_correct'] == False]['user_id'].nunique(),
            'Percent Correct': x['is_correct'].mean() * 100,
            'Modal Answer': x['user_answer'].mode()[0] if not x['user_answer'].mode().empty else None
        }), include_groups=False).reset_index()

        # Save summary statistics to a CSV file
        summary_output_path = output_dir / 'summary_statistics.csv'
        summary.to_csv(summary_output_path, index=False)

        # Generate the HTML report
        generate_html_report(df, summary, output_dir)
        # Generate the dashboard
        generate_dashboard(df, summary)


# %%
if __name__ == "__main__":

    from langchain_community.llms import Ollama
    from langchain_openai import ChatOpenAI
    from pyprojroot.here import here
    wd = here()

    llm = ChatOpenAI(
        model="gpt-4o-mini",
        temperature=0.8,
        max_tokens=None,
        timeout=None,
        max_retries=2,
        api_key=OPENAI_KEY,
        organization=OPENAI_ORG,
    )
    # llm = Ollama(
    #     model="llama3.1",

    #     temperature=0.3
    # )
    lesson_no = 20

    # Path definitions
    readingDir = Path(os.getenv('readingsDir'))
    slideDir = Path(os.getenv('slideDir'))
    syllabus_path = Path(os.getenv('syllabus_path'))

    maker = QuizMaker(llm=llm,
                      syllabus_path=syllabus_path,
                      reading_dir=readingDir,
                      output_dir=wd/"ClassFactoryOutput/QuizMaker",
                      quiz_prompt=quiz_prompt,
                      lesson_range=range(19, 21),
                      course_name="American Government",
                      prior_quiz_path=outputDir)

    # quiz_name = wd / f"ClassFactoryOutput/QuizMaker/L24/l22_24_quiz.xlsx"
    # quiz_path = wd / f"ClassFactoryOutput/QuizMaker/"

    quiz = maker.make_a_quiz()
    # maker.save_quiz_to_ppt(quiz)

    # maker.save_quiz(quiz)
    # maker.launch_interactive_quiz(quiz_name, sample_size=5, save_results=True,
    #                               output_dir=quiz_path, qr_name="test_quiz")
    # maker.assess_quiz_results()
